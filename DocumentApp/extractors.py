import os
import logging
import pandas as pd
import pdfplumber
import docx
from pptx import Presentation
from PIL import Image
import easyocr
import numpy as np
from pdf2image import convert_from_path

logger = logging.getLogger(__name__)

reader = easyocr.Reader(['en'], gpu=True)


def extract_text_from_pdf(path):
    texts = []
    tables = []

    try:
        with pdfplumber.open(path) as pdf:
            for page in pdf.pages:
                text = page.extract_text()
                if text:
                    texts.append(text)

                try:
                    page_tables = page.extract_tables() or []
                    if page_tables:
                        tables.extend(page_tables)
                except Exception:
                    logger.exception("Table extraction failed on page")

    except Exception:
        logger.exception("pdfplumber failed to open PDF")

    combined_text = "\n".join(texts).strip()

    if not combined_text:
        logger.info("No embedded text found. Running OCR on PDF pages...")

        try:
            images = convert_from_path(path, dpi=300)
            ocr_texts = []

            for img in images:
                img_np = np.array(img)
                ocr_result = reader.readtext(img_np, detail=0, paragraph=True)
                ocr_texts.append("\n".join(ocr_result))

            combined_text = "\n".join(ocr_texts)

        except Exception:
            logger.exception("OCR fallback failed on PDF")
            combined_text = ""

    return combined_text, {"tables_extracted": len(tables)}


def extract_text_from_docx(path):
    try:
        doc = docx.Document(path)
        paras = [p.text for p in doc.paragraphs if p.text.strip()]
        return "\n\n".join(paras), {}
    except Exception:
        logger.exception("DOCX extraction failed")
        return "", {}


def extract_text_from_pptx(path):
    try:
        prs = Presentation(path)
        slides_text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slides_text.append(shape.text)
        return "\n\n".join(slides_text), {}
    except Exception:
        logger.exception("PPTX extraction failed")
        return "", {}


def extract_text_from_csv(path):
    try:
        df = pd.read_csv(path)
        text = df.astype(str).fillna("").to_string(index=False)
        return text, {"columns": list(df.columns)}
    except Exception:
        logger.exception("CSV extraction failed")
        return "", {}


def extract_text_from_image(path):
    try:
        result = reader.readtext(path, detail=0, paragraph=True)
        text = "\n".join(result)
        return text, {}
    except Exception:
        logger.exception("Image OCR failed")
        return "", {}


def extract_text_from_file(path_or_file):
    if hasattr(path_or_file, "path"):
        path = path_or_file.path
    elif isinstance(path_or_file, str):
        path = path_or_file
    else:
        path = None

    if not path or not os.path.exists(path):
        logger.error(f"File path not found: {path}")
        return "", {}

    ext = path.split(".")[-1].lower()

    if ext == "pdf":
        return extract_text_from_pdf(path)

    if ext == "docx":
        return extract_text_from_docx(path)

    if ext in ("ppt", "pptx"):
        return extract_text_from_pptx(path)

    if ext == "csv":
        return extract_text_from_csv(path)

    if ext in ("png", "jpg", "jpeg", "bmp", "tiff"):
        return extract_text_from_image(path)

    try:
        with open(path, "r", encoding="utf-8", errors="ignore") as f:
            return f.read(), {}
    except Exception:
        return "", {}
